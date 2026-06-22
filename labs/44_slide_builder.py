"""Document & Slide Generation — Session 44: Slide Deck Builder.

Two-stage pipeline:
  Stage 1 — Claude generates a structured JSON outline (titles, bullets, speaker notes)
  Stage 2 — python-pptx renders the outline into a .pptx file

Usage:
    pip install python-pptx
    python 44_slide_builder.py --prompt "10-slide deck on RAG architecture for engineers"
    python 44_slide_builder.py --prompt "..." --slides 8 --output deck.pptx
"""

import argparse
import json
from pathlib import Path

from dotenv import load_dotenv
from langchain_anthropic import ChatAnthropic
from langchain_core.messages import HumanMessage, SystemMessage
from pydantic import BaseModel, Field

load_dotenv()

MODEL = "claude-opus-4-7"
llm = ChatAnthropic(model=MODEL, temperature=0, max_tokens=4096)


class Slide(BaseModel):
    title: str = Field(description="Slide title (short, max 8 words)")
    bullets: list[str] = Field(description="3-5 bullet points (concise, action-oriented)")
    speaker_notes: str = Field(description="2-4 sentences of speaker notes for this slide")


class DeckOutline(BaseModel):
    deck_title: str
    slides: list[Slide]


OUTLINE_SYSTEM = """You are a presentation designer. You create clear, logical slide outlines.

Rules for bullet points:
- Start each with an action verb or key noun
- Max 10 words per bullet
- No full sentences — fragments are better on slides
- 3-5 bullets per slide

Rules for speaker notes:
- Written for the speaker, not the audience
- Include context, examples, or transitions
- 2-4 sentences"""

OUTLINE_PROMPT = """Create a {num_slides}-slide deck for this request:

{prompt}

Return a JSON object with:
{{
  "deck_title": "The Deck Title",
  "slides": [
    {{
      "title": "Slide Title",
      "bullets": ["Bullet 1", "Bullet 2", "Bullet 3"],
      "speaker_notes": "Speaker notes for this slide..."
    }},
    ...
  ]
}}

Include exactly {num_slides} slides."""


def generate_outline(prompt: str, num_slides: int) -> DeckOutline:
    response = llm.invoke([
        SystemMessage(content=OUTLINE_SYSTEM),
        HumanMessage(content=OUTLINE_PROMPT.format(prompt=prompt, num_slides=num_slides)),
    ])
    raw = response.content
    start = raw.find("{")
    end = raw.rfind("}") + 1
    data = json.loads(raw[start:end])
    return DeckOutline(**data)


def build_pptx(outline: DeckOutline, output_path: str) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = outline.deck_title
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"{len(outline.slides)} slides"

    # Content slides
    content_layout = prs.slide_layouts[1]
    for slide_data in outline.slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data.title

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for i, bullet in enumerate(slide_data.bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
            tf.paragraphs[i].level = 0

        # Speaker notes
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = slide_data.speaker_notes

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate a slide deck from a prompt.")
    parser.add_argument("--prompt", required=True, help="What the deck should cover")
    parser.add_argument("--slides", type=int, default=10, help="Number of slides (default: 10)")
    parser.add_argument("--output", default="deck.pptx", help="Output .pptx file")
    parser.add_argument("--outline-only", action="store_true", help="Print outline JSON without building PPTX")
    args = parser.parse_args()

    print(f"Generating {args.slides}-slide outline...")
    outline = generate_outline(args.prompt, args.slides)

    print(f"\nDeck: {outline.deck_title}")
    for i, slide in enumerate(outline.slides, 1):
        print(f"  {i}. {slide.title}")

    if args.outline_only:
        print("\n" + json.dumps(outline.model_dump(), indent=2))
        return

    print(f"\nBuilding {args.output}...")
    build_pptx(outline, args.output)
    print(f"Deck written to {args.output} ({len(outline.slides)} content slides + title slide)")


if __name__ == "__main__":
    main()
